"""
Padronizador de Documentos — Amélie & Juliette
Backend FastAPI v9 — HTML quality output + editável + logos oficiais + imagens
"""

import os, io, json, re, base64, zipfile
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
import anthropic
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
import openpyxl
from PIL import Image as PILImage

app = FastAPI(title="Padronizador Amélie & Juliette v9")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))

from logos import LOGO_AMELIE, LOGO_JULIETTE

BRAND = {
    "amelie": {
        "name": "Amélie", "logo_b64": LOGO_AMELIE,
        "primary": "#C8102E", "primary_rgb": (200,16,46),
        "text": "#1a1a1a", "muted": "#999999",
        "bg": "#ffffff", "bg2": "#fafafa", "border": "#C8102E",
        "footer_bg": "#fafafa", "footer_text": "#999999",
        "footer_logo_filter": "none", "accent_hex": "C8102E",
        "persona": """Você representa a marca Amélie.
IDENTIDADE: sofisticada, acolhedora, gastronômica. Linguagem elegante sem ser distante.
PADRÃO VISUAL: vermelho #C8102E, tipografia sem serifa em caixa alta, tabelas com header vermelho, passos numerados com círculos vermelhos."""
    },
    "juliette": {
        "name": "Juliette", "logo_b64": LOGO_JULIETTE,
        "primary": "#506775", "primary_rgb": (80,103,117),
        "text": "#506775", "muted": "#BFA56E",
        "bg": "#F5F6F6", "bg2": "#ffffff", "border": "#BFA56E",
        "footer_bg": "#506775", "footer_text": "#BFA56E",
        "footer_logo_filter": "brightness(0) invert(1)", "accent_hex": "BFA56E",
        "persona": """Você representa a marca Juliette Bistrô Art Déco.
IDENTIDADE: nostalgia, sofisticação, glamour anos 60. Linguagem como poesia urbana — calma, sedutora, elegante.
PADRÃO VISUAL: azul #506775, dourado #BFA56E, fundo off-white #F5F6F6, tipografia leve e itálica, separadores dourados."""
    }
}

AUD = {
    "franqueados": {"label": "Franqueados", "instrucoes": "PÚBLICO: Franqueados. Entregue apenas o que precisam saber para operar. Linguagem contratualmente segura: 'previsto em contrato', 'conforme manual operacional'."},
    "time": {"label": "Time interno", "instrucoes": "PÚBLICO: Time interno. Quadro completo: contexto, impactos, dependências, responsáveis, prazos. Visão 360 sem omitir nada."},
    "liderancas": {"label": "Lideranças", "instrucoes": "PÚBLICO: Lideranças. Pareto: o que mais impacta aparece primeiro. Headline → diagnóstico → riscos → ações → decisões pendentes."}
}

DOC_LABELS = {"ficha": "ficha técnica", "relatorio": "relatório", "checklist": "checklist", "apresentacao": "apresentação"}

INTERVENCAO = """ESCOPO — REGRA FUNDAMENTAL:
NÃO altere, resuma, expanda ou reescreva o conteúdo. Preserve TUDO integralmente.
Atue apenas em: formatação da marca, organização estrutural, ortografia e gramática."""

# ── Extratores ───────────────────────────────────────────────────────────────

def extract_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    lines, images, img_counter, seen = [], [], [0], set()
    for para in doc.paragraphs:
        style, text = para.style.name, para.text.strip()
        if not text: lines.append(""); continue
        if "Heading 1" in style: lines.append(f"# {text}")
        elif "Heading 2" in style: lines.append(f"## {text}")
        elif "Heading 3" in style: lines.append(f"### {text}")
        elif style.startswith("List"): lines.append(f"- {text}")
        else: lines.append(text)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                key = hash(img_data)
                if key not in seen:
                    seen.add(key)
                    b64 = base64.b64encode(img_data).decode()
                    ext = rel.target_ref.split('.')[-1].lower()
                    images.append({'b64': b64, 'ext': ext, 'id': img_counter[0]})
                    img_counter[0] += 1
            except: pass
    return "\n".join(lines), images

def is_slide_pdf(file_bytes):
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            if not pdf.pages: return False
            page = pdf.pages[0]
            is_land = (page.width or 0) > (page.height or 0)
            total = sum(len((p.extract_text() or "").split()) for p in pdf.pages)
            avg = total / len(pdf.pages) if pdf.pages else 0
            return is_land and avg < 120
    except: return False

def extract_pdf(file_bytes):
    text_parts, images, img_counter = [], [], 0
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            slide_mode = is_slide_pdf(file_bytes)
            for i, page in enumerate(pdf.pages, 1):
                t = page.extract_text()
                if t and t.strip():
                    prefix = f"## Slide {i}\n\n" if slide_mode else ""
                    text_parts.append(f"{prefix}{t.strip()}")
                for img in page.images:
                    try:
                        x0,y0,x1,y1 = img['x0'],img['y0'],img['x1'],img['y1']
                        if (x1-x0)>50 and (y1-y0)>50:
                            cropped = page.crop((x0,y0,x1,y1)).to_image(resolution=150)
                            buf = io.BytesIO()
                            cropped.save(buf, format='PNG')
                            b64 = base64.b64encode(buf.getvalue()).decode()
                            images.append({'b64': b64, 'ext': 'png', 'id': img_counter})
                            img_counter += 1
                    except: pass
    except Exception as e: raise HTTPException(422, f"Erro PDF: {e}")
    return "\n\n".join(text_parts), images

def extract_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    lines, images, img_counter = [], [], 0
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape,"text") and shape.text.strip(): lines.append(shape.text.strip())
            if shape.shape_type == 13:
                try:
                    b64 = base64.b64encode(shape.image.blob).decode()
                    images.append({'b64': b64, 'ext': shape.image.ext, 'id': img_counter})
                    img_counter += 1
                except: pass
    return "\n\n".join(lines), images

def extract_xlsx(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"## Aba: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            if any(v.strip() for v in vals): lines.append(" | ".join(vals))
    return "\n".join(lines), []

# ── Claude ───────────────────────────────────────────────────────────────────

def padronizar_claude(text, brand, audience, doc_type, reference="", images=[]):
    b, a = BRAND[brand], AUD[audience]
    ref_block = f"\n\nREFERÊNCIA:\n---\n{reference}\n---" if reference.strip() else ""
    img_note = f"\n\nO documento contém {len(images)} imagem(ns). Preserve [IMAGEM_N] onde estavam." if images else ""
    prompt = f"""Você é especialista em padronização de documentos corporativos.
MARCA: {b['name']}
{b['persona']}
{a['instrucoes']}
TIPO: {DOC_LABELS.get(doc_type, doc_type)}
{INTERVENCAO}{ref_block}{img_note}
DOCUMENTO:
---
{text}
---
Responda SOMENTE com JSON válido:
{{"documento_padronizado":"markdown completo","tipo_detectado":"documento ou slides","titulo":"título principal","subtitulo":"subtítulo ou categoria","tag":"categoria curta","rodape":"texto de rodapé","alteracoes":"lista de alterações (máx 5, cada uma com '-')"}}"""
    resp = client.messages.create(model="claude-sonnet-4-5", max_tokens=4000, messages=[{"role":"user","content":prompt}])
    raw = resp.content[0].text.replace("```json","").replace("```","").strip()
    return json.loads(raw)

# ── Render markdown ──────────────────────────────────────────────────────────

def render_md(markdown, b, img_map):
    lines = markdown.split("\n")
    html, in_list, in_table, table_rows = [], False, False, []

    def close_list():
        nonlocal in_list
        if in_list: html.append("</ul>"); in_list = False

    def close_table():
        nonlocal in_table, table_rows
        if not in_table or not table_rows: return
        t = f'<table style="width:100%;border-collapse:collapse;margin:16px 0;font-size:13px">'
        for ri, row in enumerate(table_rows):
            cells = [c.strip() for c in row.split("|") if c.strip()]
            if not cells: continue
            if ri == 0:
                t += "<thead><tr>" + "".join(f'<th style="background:{b["primary"]};color:#fff;font-weight:600;font-size:10px;letter-spacing:1px;text-transform:uppercase;padding:10px 14px;text-align:left">{c}</th>' for c in cells) + "</tr></thead><tbody>"
            elif ri == 1 and all(set(c.strip()) <= set('-:|') for c in cells): continue
            else: t += "<tr>" + "".join(f'<td style="padding:9px 14px;border-bottom:1px solid #f0f0f0;color:#444">{c}</td>' for c in cells) + "</tr>"
        t += "</tbody></table>"
        html.append(t); table_rows.clear(); in_table = False

    for line in lines:
        img_m = re.match(r'\[IMAGEM_(\d+)\]', line.strip())
        if img_m:
            close_list(); close_table()
            img_id = int(img_m.group(1))
            if img_id in img_map:
                img = img_map[img_id]
                html.append(f'<div style="margin:16px 0;text-align:center"><img src="data:image/{img["ext"]};base64,{img["b64"]}" style="max-width:100%;height:auto;border-radius:4px"></div>')
            continue
        if "|" in line and line.strip().startswith("|"):
            close_list(); in_table = True; table_rows.append(line.strip().strip("|")); continue
        else: close_table()

        if line.startswith("### "):
            close_list(); html.append(f'<h3 style="font-size:11px;font-weight:700;letter-spacing:2px;text-transform:uppercase;color:{b["primary"]};margin:20px 0 8px">{line[4:]}</h3>')
        elif line.startswith("## "):
            close_list(); html.append(f'<h2 style="font-size:15px;font-weight:600;color:{b["text"]};margin:24px 0 8px">{line[3:]}</h2>')
        elif line.startswith("# "):
            close_list(); html.append(f'<h1 style="font-size:22px;font-weight:700;color:{b["primary"]};margin:0 0 12px;text-align:center">{line[2:]}</h1>')
        elif line.startswith("- ") or line.startswith("* "):
            if not in_list: html.append('<ul style="margin:8px 0 12px;padding:0;list-style:none">'); in_list = True
            text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', line[2:])
            if b["primary"] == "#C8102E":
                html.append(f'<li style="display:flex;gap:10px;margin-bottom:7px;align-items:flex-start"><span style="width:6px;height:6px;background:{b["primary"]};border-radius:50%;flex-shrink:0;margin-top:6px"></span><span style="font-size:13px;color:#444;line-height:1.5">{text}</span></li>')
            else:
                html.append(f'<li style="font-size:13px;color:{b["text"]};padding:3px 0 3px 12px;border-left:2px solid {b["border"]};margin-bottom:6px">{text}</li>')
        elif re.match(r'^\d+\.', line.strip()):
            close_list()
            num = re.match(r'^(\d+)\.', line.strip()).group(1)
            text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', re.sub(r'^\d+\.\s*','',line.strip()))
            html.append(f'<div style="display:flex;gap:12px;margin-bottom:10px;align-items:flex-start"><div style="width:24px;height:24px;background:{b["primary"]};color:#fff;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:11px;font-weight:700;flex-shrink:0">{num}</div><div style="font-size:13px;color:#444;line-height:1.5;padding-top:4px">{text}</div></div>')
        elif line.strip() == "---":
            close_list(); html.append(f'<hr style="border:none;border-top:1px solid {b["border"]};opacity:0.3;margin:20px 0">')
        elif not line.strip():
            close_list(); html.append('<div style="height:6px"></div>')
        else:
            close_list()
            text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', line)
            if text.strip(): html.append(f'<p style="font-size:13px;color:#444;line-height:1.6;margin:0 0 8px">{text}</p>')

    close_list(); close_table()
    return "\n".join(html)

# ── Build HTML ────────────────────────────────────────────────────────────────

def build_html(parsed, brand, images, is_slides):
    b = BRAND[brand]
    img_map = {img['id']: img for img in images}
    logo = f'<img src="data:image/png;base64,{b["logo_b64"]}" style="height:38px;width:auto;display:block">'
    logo_f = f'<img src="data:image/png;base64,{b["logo_b64"]}" style="height:22px;width:auto;filter:{b["footer_logo_filter"]};opacity:0.85">'
    titulo = parsed.get("titulo",""); subtitulo = parsed.get("subtitulo","")
    tag = parsed.get("tag",""); rodape = parsed.get("rodape", b["name"]+" · 2026")
    md = parsed.get("documento_padronizado","")
    fonts = '<link href="https://fonts.googleapis.com/css2?family=Montserrat:ital,wght@0,300;0,400;0,500;0,600;0,700;1,300;1,400&display=swap" rel="stylesheet">'

    if is_slides:
        slides_data = []
        current = {"title":"","body":[]}
        for line in md.split("\n"):
            if line.startswith("## ") or line.startswith("# "):
                if current["title"] or current["body"]: slides_data.append(dict(current))
                current = {"title": line.lstrip("# "), "body":[]}
            else: current["body"].append(line)
        if current["title"] or current["body"]: slides_data.append(current)

        slides_html = ""
        for s in slides_data:
            body_html = render_md("\n".join(s["body"]), b, img_map)
            slides_html += f"""<div class="slide">
  <div class="slide-header"><div>{logo}</div><div class="tag">{tag}</div></div>
  <div class="slide-body"><h2 class="slide-title">{s["title"]}</h2>{body_html}</div>
  <div class="slide-footer">{logo_f}<div class="footer-text">{rodape}</div></div>
</div>"""

        return f"""<!DOCTYPE html><html lang="pt-BR"><head><meta charset="UTF-8"><title>{titulo or b["name"]}</title>{fonts}
<style>*{{margin:0;padding:0;box-sizing:border-box}}body{{font-family:'Montserrat',sans-serif;background:#e0e0de;padding:1rem}}
.slide{{background:{b["bg"]};margin-bottom:3px;display:flex;flex-direction:column;min-height:520px;max-width:960px;margin-left:auto;margin-right:auto;margin-bottom:16px;box-shadow:0 2px 20px rgba(0,0,0,0.1)}}
.slide-header{{padding:16px 48px 12px;border-bottom:2px solid {b["primary"]};display:flex;justify-content:space-between;align-items:center;background:{b["bg2"]}}}
.tag{{font-size:9px;color:{b["muted"]};letter-spacing:2px;text-transform:uppercase}}
.slide-body{{flex:1;padding:40px 80px;display:flex;flex-direction:column;justify-content:center}}
.slide-title{{font-size:30px;font-weight:700;color:{b["primary"]};margin-bottom:24px;text-align:center}}
.slide-footer{{padding:12px 48px;border-top:2px solid {b["primary"]};display:flex;justify-content:space-between;align-items:center;background:{b["footer_bg"]}}}
.footer-text{{font-size:9px;letter-spacing:2px;text-transform:uppercase;color:{b["footer_text"]}}}
@media print{{body{{background:white;padding:0}}.slide{{box-shadow:none;page-break-after:always}}}}
</style></head><body>{slides_html}</body></html>"""

    else:
        tag_html = f'<div style="display:inline-block;font-size:9px;font-weight:700;letter-spacing:3px;text-transform:uppercase;color:{b["primary"]};border:1px solid {b["primary"]};padding:3px 10px;margin-bottom:18px">{tag}</div>' if tag else ""
        titulo_html = f'<div style="font-size:26px;font-weight:700;color:{b["text"]};margin-bottom:4px">{titulo}</div>' if titulo else ""
        sub_html = f'<div style="font-size:11px;font-weight:500;letter-spacing:2px;text-transform:uppercase;color:{b["muted"]};margin-bottom:24px">{subtitulo}</div>' if subtitulo else ""
        content = render_md(md, b, img_map)

        return f"""<!DOCTYPE html><html lang="pt-BR"><head><meta charset="UTF-8"><title>{titulo or b["name"]}</title>{fonts}
<style>*{{margin:0;padding:0;box-sizing:border-box}}body{{font-family:'Montserrat',sans-serif;background:#e0e0de;padding:2rem 1rem}}
.page{{background:{b["bg"]};max-width:800px;margin:0 auto;box-shadow:0 2px 24px rgba(0,0,0,0.1)}}
@media print{{body{{background:white;padding:0}}.page{{box-shadow:none;max-width:100%}}}}
</style></head><body><div class="page">
  <div style="padding:20px 48px 16px;border-bottom:2px solid {b["primary"]};display:flex;justify-content:space-between;align-items:center;background:{b["bg2"]}">{logo}
    <div style="font-size:10px;color:{b["muted"]};letter-spacing:2px;text-transform:uppercase">2026</div></div>
  <div style="padding:32px 48px;background:{b["bg"]}">{tag_html}{titulo_html}{sub_html}{content}</div>
  <div style="padding:16px 48px;border-top:2px solid {b["primary"]};display:flex;justify-content:space-between;align-items:center;background:{b["footer_bg"]}">{logo_f}
    <div style="font-size:9px;letter-spacing:2px;text-transform:uppercase;color:{b["footer_text"]}">{rodape}</div></div>
</div></body></html>"""

# ── Build DOCX ────────────────────────────────────────────────────────────────

def build_docx_edit(markdown, brand, images):
    b = BRAND[brand]; r,g,bl = b["primary_rgb"]
    doc = Document()
    style = doc.styles["Normal"]; style.font.name="Calibri"; style.font.size=Pt(11)
    for h,size,bold in [("Heading 1",18,True),("Heading 2",14,True),("Heading 3",12,False)]:
        hs=doc.styles[h]; hs.font.name="Calibri"; hs.font.size=Pt(size); hs.font.bold=bold; hs.font.color.rgb=RGBColor(r,g,bl)
    try:
        section=doc.sections[0]; header=section.header
        hp=header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        hp.clear(); run=hp.add_run()
        run.add_picture(io.BytesIO(base64.b64decode(b["logo_b64"])), height=Cm(1.0))
        hp.alignment=WD_ALIGN_PARAGRAPH.LEFT
        pPr=hp._p.get_or_add_pPr(); pBdr=OxmlElement('w:pBdr')
        bottom=OxmlElement('w:bottom'); bottom.set(qn('w:val'),'single'); bottom.set(qn('w:sz'),'6'); bottom.set(qn('w:color'),b["accent_hex"])
        pBdr.append(bottom); pPr.append(pBdr)
    except: pass
    img_map = {img['id']: img for img in images}
    for line in markdown.split("\n"):
        img_m = re.match(r'\[IMAGEM_(\d+)\]', line.strip())
        if img_m:
            img_id=int(img_m.group(1))
            if img_id in img_map:
                try:
                    p=doc.add_paragraph(); run=p.add_run()
                    run.add_picture(io.BytesIO(base64.b64decode(img_map[img_id]['b64'])), width=Cm(12))
                    p.alignment=WD_ALIGN_PARAGRAPH.CENTER
                except: pass
            continue
        line=line.rstrip()
        if line.startswith("### "): doc.add_heading(line[4:],level=3)
        elif line.startswith("## "): doc.add_heading(line[3:],level=2)
        elif line.startswith("# "): doc.add_heading(line[2:],level=1)
        elif line.startswith("- ") or line.startswith("* "): doc.add_paragraph(line[2:],style="List Bullet")
        elif re.match(r"^\d+\. ",line): doc.add_paragraph(re.sub(r"^\d+\. ","",line),style="List Number")
        elif not line.strip(): doc.add_paragraph()
        else:
            p=doc.add_paragraph(); parts=re.split(r"\*\*(.*?)\*\*",line)
            for i,part in enumerate(parts): run=p.add_run(part); run.bold=(i%2==1)
    buf=io.BytesIO(); doc.save(buf); buf.seek(0); return buf.read()

# ── Build PPTX ────────────────────────────────────────────────────────────────

def build_pptx_edit(markdown, brand, images):
    b=BRAND[brand]; r,g,bl=b["primary_rgb"]
    prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    logo_bytes=base64.b64decode(b["logo_b64"])
    slides_data=[]; current={"title":"","body":[]}
    for line in markdown.split("\n"):
        if line.startswith("## ") or line.startswith("# "):
            if current["title"] or current["body"]: slides_data.append(dict(current))
            current={"title":line.lstrip("# "),"body":[]}
        else: current["body"].append(line)
    if current["title"] or current["body"]: slides_data.append(current)
    for sd in slides_data:
        slide=prs.slides.add_slide(prs.slide_layouts[6])
        try: slide.shapes.add_picture(io.BytesIO(logo_bytes),Inches(0.3),Inches(0.15),height=Inches(0.5))
        except: pass
        ls=slide.shapes.add_shape(1,Inches(0),Inches(0.75),Inches(13.33),Inches(0.02))
        ls.fill.solid(); ls.fill.fore_color.rgb=PptRGB(r,g,bl); ls.line.fill.background()
        tb=slide.shapes.add_textbox(Inches(1),Inches(1.2),Inches(11.33),Inches(1.2))
        p=tb.text_frame.paragraphs[0]; p.text=sd["title"]
        p.font.size=PptPt(32); p.font.bold=True; p.font.color.rgb=PptRGB(r,g,bl); p.alignment=1
        body_text="\n".join([l.lstrip("- *").lstrip("0123456789. ") for l in sd["body"] if l.strip() and not re.match(r'\[IMAGEM_\d+\]',l.strip())])
        if body_text:
            bb=slide.shapes.add_textbox(Inches(1),Inches(2.6),Inches(11.33),Inches(4.0))
            tf=bb.text_frame; tf.word_wrap=True; tf.paragraphs[0].text=body_text
            tf.paragraphs[0].font.size=PptPt(18); tf.paragraphs[0].font.color.rgb=PptRGB(60,60,60)
    buf=io.BytesIO(); prs.save(buf); buf.seek(0); return buf.read()

# ── Frontend ──────────────────────────────────────────────────────────────────

_static = os.path.join(os.path.dirname(__file__), "static", "index.html")
FRONTEND_HTML = open(_static).read() if os.path.exists(_static) else "<h1>Frontend não encontrado</h1>"

@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    return HTMLResponse(content=FRONTEND_HTML)

# ── Endpoint principal ────────────────────────────────────────────────────────

@app.post("/padronizar")
async def padronizar(
    file: UploadFile = File(...),
    brand: str = Form(...),
    audience: str = Form(...),
    doc_type: str = Form(...),
    reference: str = Form(""),
):
    file_bytes = await file.read()
    filename = file.filename or "documento"
    ext = filename.rsplit(".",1)[-1].lower()
    if brand not in BRAND: raise HTTPException(400,"Marca inválida")
    if audience not in AUD: raise HTTPException(400,"Público inválido")

    # Estimativa de tempo
    size_kb = len(file_bytes)/1024
    est = 20 + (10 if size_kb > 500 else 0) + (5 if ext in ['pdf','pptx'] else 0)

    # Extrai
    try:
        if ext=="docx": text,images=extract_docx(file_bytes)
        elif ext=="pdf": text,images=extract_pdf(file_bytes)
        elif ext=="pptx": text,images=extract_pptx(file_bytes)
        elif ext in ("xlsx","xls"): text,images=extract_xlsx(file_bytes)
        elif ext=="txt": text,images=file_bytes.decode("utf-8",errors="replace"),[]
        else: raise HTTPException(400,f"Formato .{ext} não suportado")
    except HTTPException: raise
    except Exception as e: raise HTTPException(422,f"Erro ao ler: {e}")

    if not text.strip(): raise HTTPException(422,"Não foi possível extrair texto")

    is_slides=(ext=="pptx") or (ext=="pdf" and is_slide_pdf(file_bytes))

    # Claude
    try: parsed=padronizar_claude(text,brand,audience,doc_type,reference,images)
    except Exception as e: raise HTTPException(500,f"Erro Claude: {e}")

    md_out=parsed.get("documento_padronizado",text)

    # HTML
    html_out=build_html(parsed,brand,images,is_slides)

    # Editável
    try:
        if is_slides: editable=build_pptx_edit(md_out,brand,images); edit_ext="pptx"
        else: editable=build_docx_edit(md_out,brand,images); edit_ext="docx"
    except: editable=None; edit_ext="docx"

    # ZIP
    base_name=filename.rsplit(".",1)[0]
    zip_buf=io.BytesIO()
    with zipfile.ZipFile(zip_buf,'w',zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{base_name}_padronizado.html", html_out.encode("utf-8"))
        if editable: zf.writestr(f"{base_name}_padronizado.{edit_ext}", editable)
        nota=f"""Documento padronizado — {BRAND[brand]['name']}
Público: {AUD[audience]['label']} | Tipo: {DOC_LABELS.get(doc_type,doc_type)}

Arquivos:
- {base_name}_padronizado.html → versão final (abra no navegador e imprima como PDF)
- {base_name}_padronizado.{edit_ext or 'docx'} → versão editável

O que foi ajustado:
{parsed.get('alteracoes','')}
"""
        zf.writestr("LEIA-ME.txt", nota.encode("utf-8"))
    zip_buf.seek(0)

    headers={
        "Content-Disposition": f'attachment; filename="{base_name}_padronizado.zip"',
        "X-Alteracoes": parsed.get("alteracoes","").replace("\n"," | ")[:400],
        "X-Tempo-Estimado": str(est),
        "X-Titulo": parsed.get("titulo","")[:100],
    }
    return StreamingResponse(zip_buf, media_type="application/zip", headers=headers)

@app.get("/health")
def health():
    return {"status":"ok","version":"v9"}
